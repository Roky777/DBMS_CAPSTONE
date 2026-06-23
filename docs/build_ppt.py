#!/usr/bin/env python3
# Generates docs/Athenaeum_Presentation.pptx — Athenaeum theme, 4-part flow,
# with real screenshots of the running app.
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SHOTS = "/tmp/slides"

CREAM = RGBColor(0xF4,0xEC,0xE0); PAPER = RGBColor(0xFB,0xF6,0xEE)
NAVY  = RGBColor(0x26,0x41,0x5E); BLUE  = RGBColor(0x2F,0x51,0x70)
TERRA = RGBColor(0xA8,0x50,0x3A); INK   = RGBColor(0x2C,0x2A,0x26)
MUTED = RGBColor(0x7A,0x72,0x63); LINE  = RGBColor(0xE6,0xDC,0xCC)
WHITE = RGBColor(0xFF,0xFF,0xFF); SOFT  = RGBColor(0xF6,0xEF,0xE4)
SAND  = RGBColor(0xD9,0xC7,0xB0)
SERIF = "Georgia"; SANS = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

PG = [0]
def slide(bg=PAPER):
    PG[0] += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s,x,y,w,h): return s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08):
    sp = s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit=False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except: pass
    return sp

def text(tb, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.04):
    tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    first=True
    for item in runs:
        txt,size,color,bold,font = (item+(None,)*5)[:5]
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.line_spacing=space; p.space_after=Pt(4)
        r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bool(bold)
        r.font.color.rgb=color or INK; r.font.name=font or SANS
    return tb

def bullets(tb, items, size=16, color=INK, gap=8, mark="—", mark_color=TERRA):
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    first=True
    for it in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.line_spacing=1.06; p.space_after=Pt(gap)
        rm=p.add_run(); rm.text=mark+"  "; rm.font.size=Pt(size); rm.font.bold=True
        rm.font.color.rgb=mark_color; rm.font.name=SANS
        if isinstance(it,tuple):
            lead,rest=it
            r1=p.add_run(); r1.text=lead; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=color; r1.font.name=SANS
            r2=p.add_run(); r2.text=rest; r2.font.size=Pt(size); r2.font.color.rgb=color; r2.font.name=SANS
        else:
            r=p.add_run(); r.text=it; r.font.size=Pt(size); r.font.color.rgb=color; r.font.name=SANS
    return tb

def kicker(s, part, label, x=0.9, y=0.62):
    tb=box(s,x,y,11,0.4); tf=tb.text_frame; tf.word_wrap=False
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    p=tf.paragraphs[0]
    if part:
        r=p.add_run(); r.text=part+"   "; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=TERRA; r.font.name=SANS
    r=p.add_run(); r.text=label.upper(); r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=MUTED; r.font.name=SANS

def title(s, txt, x=0.9, y=0.95, w=11.5, size=34, color=NAVY):
    text(box(s,x,y,w,1.0), [(txt,size,color,True,SERIF)])

def bar(s, x=0.92, y=1.72, w=0.9):
    rect(s,x,y,w,0.06,fill=TERRA,shape=MSO_SHAPE.RECTANGLE)

def footer(s,n=None):
    text(box(s,0.9,7.04,8,0.33),[("Athenaeum · Library Management System",10,MUTED,False,SANS)])
    text(box(s,11.6,7.04,1.0,0.33),[(str(PG[0]),10,MUTED,False,SANS)],align=PP_ALIGN.RIGHT)

def shot(s, name, x, y, w):
    """place a screenshot (1380x900 → ratio 1.5333) with a soft frame."""
    h = w/1.5333
    rect(s, x-0.06, y-0.06, w+0.12, h+0.12, fill=WHITE, line=LINE, lw=1.2, radius=0.03)
    p = os.path.join(SHOTS, name)
    s.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    return h

# ============================== 1 COVER
s = slide(NAVY)
rect(s,0,4.7,13.333,2.8,fill=BLUE,shape=MSO_SHAPE.RECTANGLE)
text(box(s,1.0,0.7,11,0.5),[("DBMS CAPSTONE PROJECT",15,SAND,True,SANS)])
text(box(s,1.0,2.2,11.3,1.6),[("Athenaeum",66,WHITE,True,SERIF)])
rect(s,1.02,3.55,1.4,0.05,fill=TERRA,shape=MSO_SHAPE.RECTANGLE)
text(box(s,1.0,3.72,11.3,0.8),[("A Library Management System built on MySQL",24,CREAM,False,SERIF)])
bullets(box(s,1.0,5.05,11,1.6),[
    ("What  ","— a normalized relational database with a live web front-end"),
    ("How  ","— MySQL · Node/Express · vanilla HTML/CSS/JS, in a 3-tier design"),
    ("Where  ","— deployed live on Vercel (app) + Aiven (database)"),
], size=16, color=CREAM, mark="›", mark_color=SAND)
text(box(s,1.0,6.9,11,0.4),[("athenaeum-roky-pauls-projects.vercel.app",13,SAND,False,SANS)])

# ============================== 2 AGENDA
s = slide(); kicker(s,"","Agenda"); title(s,"How this talk is organised"); bar(s)
parts = [
    ("PART 1","What is it","The problem, the goal, and why a library is the perfect model.", NAVY),
    ("PART 2","How we did it","Architecture, the data model, and normalization to 3NF.", TERRA),
    ("PART 3","What we built","The database, the logic inside it, the live app & deployment.", BLUE),
    ("PART 4","Walkthrough","The build, step by step — design → database → app → ship.", NAVY),
]
x=0.92
for tag,name,desc,c in parts:
    rect(s,x,2.3,2.92,3.7,fill=WHITE,line=LINE,lw=1.2,radius=0.05)
    rect(s,x,2.3,2.92,0.16,fill=c,shape=MSO_SHAPE.RECTANGLE)
    text(box(s,x+0.28,2.65,2.5,0.4),[(tag,12,TERRA,True,SANS)])
    text(box(s,x+0.28,3.05,2.45,0.9),[(name,21,NAVY,True,SERIF)])
    text(box(s,x+0.28,4.05,2.45,1.8),[(desc,13.5,MUTED,False,SANS)],space=1.12)
    x+=3.06
footer(s,2)

# ============================== 3 [P1] WHAT IS IT  (+ home shot)
s = slide(); kicker(s,"PART 1","What is it"); title(s,"A working library, run by a database"); bar(s)
text(box(s,0.92,1.95,5.4,2.0),[
    ("Athenaeum manages the everyday work of a library — and every action is "
     "stored, checked and calculated by the MySQL database itself.",16,INK,False,SANS)],space=1.15)
bullets(box(s,0.92,3.25,5.4,3.4),[
    "Catalogue books, authors, publishers & categories",
    "Track each physical copy and its availability",
    "Issue & return books (borrow–return transactions)",
    "Charge late fines automatically",
    "Register members & reserve titles",
], size=15.5, gap=11)
shot(s,"app_home.png",6.7,2.0,5.9)
text(box(s,6.7,5.95,5.9,0.4),[("The live home page — bookshelf hero + today's figures",12,MUTED,False,SANS)],align=PP_ALIGN.CENTER)
footer(s,3)

# ============================== 4 [P1] WHY A LIBRARY
s = slide(); kicker(s,"PART 1","Why a library"); title(s,"It exercises every relationship type"); bar(s)
text(box(s,0.92,1.95,11.4,0.6),[("One domain that naturally needs 1:M, M:N, recursive and 1:1 links:",17,INK,False,SANS)])
cards=[("1 : M","One-to-Many","A publisher has many books; a book has many copies.",NAVY),
       ("M : N","Many-to-Many","A book has many authors; an author writes many books.",TERRA),
       ("self","Recursive","A category can be a sub-category of another.",BLUE),
       ("1 : 1","One-to-One","Each loan has at most one fine.",NAVY)]
x=0.92
for tag,name,desc,c in cards:
    rect(s,x,2.8,2.92,3.4,fill=WHITE,line=LINE,lw=1.2,radius=0.06)
    rect(s,x,2.8,2.92,0.16,fill=c,shape=MSO_SHAPE.RECTANGLE)
    text(box(s,x+0.25,3.18,2.5,0.9),[(tag,30,c,True,SERIF)])
    text(box(s,x+0.25,4.05,2.5,0.5),[(name,16,INK,True,SANS)])
    text(box(s,x+0.25,4.6,2.45,1.5),[(desc,13.5,MUTED,False,SANS)],space=1.1)
    x+=3.06
footer(s,4)

# ============================== 5 [P2] ARCHITECTURE
s = slide(); kicker(s,"PART 2","How we did it"); title(s,"System architecture (3-tier)"); bar(s)
text(box(s,0.92,1.95,11.4,0.6),[("A browser cannot reach MySQL directly — a server sits in between.",17,INK,False,SANS)])
tiers=[("Presentation","Browser\nHTML · CSS · JS","renders dashboard,\ntables & forms",NAVY),
       ("Application","Express API\n(Node.js)","calls views &\nstored procedures",TERRA),
       ("Data","MySQL\n(library_db)","tables, views,\ntriggers, procedures",BLUE)]
x=1.15
for i,(lab,mid,sub,c) in enumerate(tiers):
    rect(s,x,3.0,3.3,2.7,fill=WHITE,line=c,lw=1.6,radius=0.06)
    text(box(s,x+0.2,3.2,2.9,0.4),[(lab.upper(),13,c,True,SANS)])
    text(box(s,x+0.2,3.7,2.9,1.0),[(mid,19,INK,True,SERIF)],space=1.0)
    text(box(s,x+0.2,4.8,2.9,0.8),[(sub,13,MUTED,False,SANS)],space=1.0)
    if i<2:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+3.32),Inches(4.1),Inches(0.55),Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb=TERRA; ar.line.fill.background(); ar.shadow.inherit=False
    x+=3.87
text(box(s,0.92,6.1,11.4,0.6),[("Hosted as: application on Vercel · database on Aiven.",15,MUTED,False,SANS)])
footer(s,5)

# ============================== [P2] ER DIAGRAM
s = slide(); kicker(s,"PART 2","How we did it"); title(s,"From ER diagram to 11 tables"); bar(s)
text(box(s,0.92,1.9,11.4,0.5),[
    ("The entity–relationship design — each box is an entity that becomes one table:",15.5,INK,False,SANS)])
# ER image (ratio 1.803:1)
erw = 9.0; erh = erw/1.803
rect(s,0.55-0.05,2.45-0.05,erw+0.1,erh+0.1,fill=WHITE,line=LINE,lw=1.2,radius=0.02)
s.shapes.add_picture(os.path.join(SHOTS,"er.png"),Inches(0.55),Inches(2.45),Inches(erw),Inches(erh))
# side note
nx = 9.85
rect(s,nx,2.45,3.0,erh,fill=SOFT,line=LINE,lw=1.2,radius=0.05)
text(box(s,nx+0.28,2.7,2.5,0.4),[("11 tables",17,TERRA,True,SANS)])
bullets(box(s,nx+0.28,3.25,2.5,3.4),[
    ("5 ","catalogue"),
    ("1 ","inventory (copies)"),
    ("2 ","people"),
    ("3 ","activity"),
], size=14, gap=9)
text(box(s,nx+0.28,5.2,2.5,1.4),[
    ("Lines show cardinality — the fork (crow's foot) marks the 'many' side.",12.5,MUTED,False,SANS)],space=1.12)
footer(s)

# ============================== 6 [P2] DATA MODEL
s = slide(); kicker(s,"PART 2","How we did it"); title(s,"The data model — 11 entities"); bar(s)
text(box(s,0.92,1.95,5.7,0.5),[("Tables, grouped by purpose",16,TERRA,True,SANS)])
bullets(box(s,0.92,2.55,5.8,4.2),[
    ("Catalogue: ","publisher, category, author, book, book_author"),
    ("Inventory: ","book_copy (each physical item + status)"),
    ("People: ","member, staff"),
    ("Activity: ","borrowing, fine, reservation"),
], size=15.5, gap=14)
rect(s,7.0,2.1,5.4,4.7,fill=SOFT,line=LINE,lw=1.2,radius=0.05)
text(box(s,7.3,2.3,4.8,0.4),[("Key relationships",15,TERRA,True,SANS)])
bullets(box(s,7.3,2.85,4.85,3.7),[
    "publisher → book  (1:M)",
    "book ↔ author  via book_author  (M:N)",
    "book → book_copy  (1:M)",
    "category → category  (recursive)",
    "borrowing → fine  (1:1)",
    "member → borrowing  (1:M)",
], size=14.5, gap=12)
footer(s,6)

# ============================== 7 [P2] NORMALIZATION
s = slide(); kicker(s,"PART 2","How we did it"); title(s,"Normalization to 3NF"); bar(s)
text(box(s,0.92,1.95,11.4,0.5),[("Each fact stored once — removing redundancy and anomalies.",16,INK,False,SANS)])
steps=[("1NF","Atomic values","Multi-valued authors & copies moved to their own tables.",NAVY),
       ("2NF","No partial dependency","Non-key columns depend on the whole composite key.",TERRA),
       ("3NF","No transitive dependency","Publisher / category details split into their own tables.",BLUE)]
x=0.92
for tag,name,desc,c in steps:
    rect(s,x,2.7,3.78,3.3,fill=WHITE,line=LINE,lw=1.2,radius=0.05)
    rect(s,x,2.7,3.78,0.16,fill=c,shape=MSO_SHAPE.RECTANGLE)
    text(box(s,x+0.3,3.05,3.2,0.7),[(tag,26,c,True,SERIF)])
    text(box(s,x+0.3,3.8,3.2,0.5),[(name,15.5,INK,True,SANS)])
    text(box(s,x+0.3,4.35,3.2,1.5),[(desc,13.5,MUTED,False,SANS)],space=1.12)
    x+=3.92
footer(s,7)

# ============================== 8 [P3] DATABASE LAYER
s = slide(); kicker(s,"PART 3","What we built"); title(s,"The database layer"); bar(s)
blocks=[("Schema","sql/01_schema.sql","11 tables · all constraints · indexes · ALTER demo"),
        ("Sample data","sql/02_data.sql","~200 meaningful rows (over the 100 minimum)"),
        ("Views","sql/03_views.sql","catalogue, availability, active loans, fines"),
        ("Queries","sql/04_queries.sql","joins · subqueries · aggregates · GROUP BY / HAVING"),
        ("Procedures","sql/05_procedures.sql","3 triggers · 3 procedures · 1 function · transactions")]
y=2.05
for name,file,desc in blocks:
    rect(s,0.92,y,11.5,0.86,fill=WHITE,line=LINE,lw=1.1,radius=0.12)
    text(box(s,1.2,y+0.17,2.3,0.6),[(name,17,NAVY,True,SERIF)])
    text(box(s,3.5,y+0.2,3.1,0.5),[(file,13,TERRA,True,SANS)])
    text(box(s,6.7,y+0.2,5.6,0.5),[(desc,13.5,MUTED,False,SANS)])
    y+=0.98
footer(s,8)

# ============================== 9 [P3] LOGIC IN DB
s = slide(); kicker(s,"PART 3","What we built"); title(s,"Logic inside the database"); bar(s)
text(box(s,0.92,1.95,11.4,0.6),[("Rules live in MySQL, so data stays correct however it is accessed.",17,INK,False,SANS)])
two=[("Triggers",["Block issuing a copy that is not Available",
                  "Auto-mark a copy 'Borrowed' when issued",
                  "Auto-free the copy + add a fine on late return"],TERRA),
     ("Procedures & transactions",["sp_issue_book / sp_return_book / sp_pay_member_fines",
                  "fn_member_balance() — outstanding dues",
                  "Wrapped in transactions with ROLLBACK (ACID)"],BLUE)]
x=0.92
for head,its,c in two:
    rect(s,x,2.75,5.65,3.6,fill=WHITE,line=LINE,lw=1.2,radius=0.05)
    rect(s,x,2.75,0.16,3.6,fill=c,shape=MSO_SHAPE.RECTANGLE)
    text(box(s,x+0.4,3.0,5.0,0.5),[(head,19,NAVY,True,SERIF)])
    bullets(box(s,x+0.4,3.7,5.0,2.5),its,size=15,gap=12)
    x+=5.95
footer(s,9)

# ============================== 10 [P3] THE APP (screenshots)
s = slide(); kicker(s,"PART 3","What we built"); title(s,"The live application"); bar(s)
shot(s,"app_catalogue.png",0.92,2.05,5.85)
text(box(s,0.92,5.98,5.85,0.4),[("Catalogue — search + live availability",12,MUTED,False,SANS)],align=PP_ALIGN.CENTER)
shot(s,"app_loans.png",6.95,2.05,5.45)
text(box(s,6.95,5.74,5.45,0.4),[("On Loan — overdue flagged, one-click return",12,MUTED,False,SANS)],align=PP_ALIGN.CENTER)
footer(s,10)

# ============================== 11 [P3] DEPLOYMENT
s = slide(); kicker(s,"PART 3","What we built"); title(s,"Hosting it online"); bar(s)
steps=[("Database → Aiven","Free cloud MySQL; imported sql/init.sql (schema + data + views + procedures)."),
       ("App → Vercel","Imported the GitHub repo; static front-end + serverless API."),
       ("Wiring","Set MYSQL_URL & DB_SSL environment variables, then redeployed."),
       ("Result","A public URL — full front-end ↔ backend ↔ live database.")]
y=2.15
for h,d in steps:
    rect(s,0.92,y,11.5,1.04,fill=WHITE,line=LINE,lw=1.1,radius=0.1)
    rect(s,0.92,y,0.16,1.04,fill=TERRA,shape=MSO_SHAPE.RECTANGLE)
    text(box(s,1.3,y+0.18,4.0,0.6),[(h,17,NAVY,True,SERIF)])
    text(box(s,5.0,y+0.22,7.2,0.7),[(d,14,MUTED,False,SANS)],space=1.06)
    y+=1.16
footer(s,11)

# ============================== 12-14 WALKTHROUGH
def wt(part_title, steps, page):
    s=slide(); kicker(s,"PART 4","Walkthrough"); title(s,part_title,size=30); bar(s)
    y=2.15
    for num,head,desc in steps:
        c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.95),Inches(y),Inches(0.55),Inches(0.55))
        c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.line.fill.background(); c.shadow.inherit=False
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(num); r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=SANS
        text(box(s,1.75,y-0.02,10.6,0.5),[(head,17,NAVY,True,SANS)])
        text(box(s,1.75,y+0.44,10.6,0.6),[(desc,13.5,MUTED,False,SANS)],space=1.05)
        y+=1.02
    footer(s,page)

wt("Step 1–4 · Plan & Design",[
    (1,"Understand the problem","List what a library does; turn the nouns into entities (tables)."),
    (2,"Map relationships","Decide 1:M, M:N, recursive and 1:1 links between entities."),
    (3,"Draw the ER diagram","Blueprint all 11 entities before writing any SQL."),
    (4,"Keys & normalization","Surrogate primary keys; design straight to 3NF."),
],12)
wt("Step 5–9 · Build the Database",[
    (5,"Create tables","CREATE TABLE with PK, FK, UNIQUE, CHECK, ENUM + indexes."),
    (6,"Insert sample data","~200 rows; a mix of returned, active & overdue loans."),
    (7,"Create views","Reusable queries for catalogue, availability, loans, fines."),
    (8,"Advanced queries","Joins, subqueries, aggregates, GROUP BY / HAVING."),
    (9,"Triggers & procedures","Move integrity + workflows into the DB, in transactions."),
],13)
wt("Step 10–13 · Build & Ship",[
    (10,"Backend API","Express endpoints that read views and call stored procedures."),
    (11,"Frontend UI","Vanilla HTML/CSS/JS — responsive, handmade-style design."),
    (12,"Connect with fetch()","JSON over REST; database errors surfaced to the user."),
    (13,"Deploy online","Database on Aiven, app on Vercel, wired via env vars."),
],14)

# ============================== 15 DELIVERABLES & TECH
s = slide(); kicker(s,"","Summary"); title(s,"Deliverables & technology"); bar(s)
text(box(s,0.92,2.1,5.7,0.4),[("Deliverables",18,TERRA,True,SANS)])
bullets(box(s,0.92,2.7,5.7,4.0),[
    "Relational schema (3NF) + data dictionary",
    "SQL: schema, data, views, queries, procedures",
    "~200-record sample dataset",
    "Working web app (frontend + API)",
    "Report & step-by-step walkthrough (PDF)",
    "This presentation",
], size=15.5, gap=10)
rect(s,7.0,2.1,5.4,4.6,fill=SOFT,line=LINE,lw=1.2,radius=0.05)
text(box(s,7.3,2.35,4.8,0.4),[("Technology",18,TERRA,True,SANS)])
bullets(box(s,7.3,2.95,4.85,3.6),[
    ("Database: ","MySQL 8 (InnoDB)"),
    ("Backend: ","Node.js + Express"),
    ("Frontend: ","HTML, CSS, JavaScript"),
    ("Hosting: ","Vercel + Aiven"),
    ("Source: ","GitHub"),
], size=15.5, gap=12)
footer(s,15)

# ============================== 16 CLOSE
s = slide(NAVY)
text(box(s,1.0,2.45,11.3,1.2),[("Thank you",54,WHITE,True,SERIF)])
rect(s,1.02,3.7,1.4,0.05,fill=TERRA,shape=MSO_SHAPE.RECTANGLE)
text(box(s,1.0,3.95,11.3,0.6),[("Athenaeum — Library Management System",20,CREAM,False,SERIF)])
text(box(s,1.0,4.75,11.3,0.5),[("Live demo:  athenaeum-roky-pauls-projects.vercel.app",15,SAND,False,SANS)])
text(box(s,1.0,5.15,11.3,0.5),[("Happy to walk through the schema, the SQL, and a live demo.",14,SAND,False,SANS)])

prs.save("Athenaeum_Presentation.pptx")
print("saved · 16 slides")
